"""多模态文档加载器。

支持的格式:
  - PDF (.pdf)  — PyPDF2 / pdfplumber
  - Word (.docx) — python-docx
  - PPT (.pptx)  — python-pptx
  - 纯文本 (.txt / .md)
  - 图片 (.png/.jpg/.bmp) — OCR 文字识别

当前为"骨架实现"：定义接口和加载逻辑，实际依赖运行时按需安装。
"""
import logging
import os
from pathlib import Path
from typing import List, Optional, Dict, Union
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class LoadedDocument:
    """加载后的文档。"""
    content: str
    metadata: dict = field(default_factory=dict)
    source_path: str = ""
    file_type: str = ""


# ========== 各格式加载器 ==========

def _load_txt(path: str) -> LoadedDocument:
    """加载纯文本文件。"""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return LoadedDocument(
        content=content,
        source_path=path,
        file_type="txt",
        metadata={"file_name": os.path.basename(path)},
    )


def _load_pdf(path: str) -> LoadedDocument:
    """加载 PDF 文件。

    优先使用 pdfplumber（表格提取好），回退 PyPDF2。
    """
    text_parts = []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
    except ImportError:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        except ImportError:
            raise ImportError(
                "PDF 加载需要 pdfplumber 或 PyPDF2。"
                "安装: pip install pdfplumber 或 pip install PyPDF2"
            )

    return LoadedDocument(
        content="\n\n".join(text_parts),
        source_path=path,
        file_type="pdf",
        metadata={
            "file_name": os.path.basename(path),
            "pages": len(text_parts),
        },
    )


def _load_docx(path: str) -> LoadedDocument:
    """加载 Word 文档。"""
    try:
        from docx import Document

        doc = Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # 提取表格
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)

        return LoadedDocument(
            content="\n\n".join(paragraphs),
            source_path=path,
            file_type="docx",
            metadata={"file_name": os.path.basename(path)},
        )
    except ImportError:
        raise ImportError(
            "Word 文档加载需要 python-docx。安装: pip install python-docx"
        )


def _load_pptx(path: str) -> LoadedDocument:
    """加载 PPT 文档。"""
    try:
        from pptx import Presentation

        prs = Presentation(path)
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())
            if slide_parts:
                slides_text.append(f"--- 第{i}页 ---\n" + "\n".join(slide_parts))

        return LoadedDocument(
            content="\n\n".join(slides_text),
            source_path=path,
            file_type="pptx",
            metadata={
                "file_name": os.path.basename(path),
                "slides": len(slides_text),
            },
        )
    except ImportError:
        raise ImportError(
            "PPT 加载需要 python-pptx。安装: pip install python-pptx"
        )


def _load_image(path: str) -> LoadedDocument:
    """加载图片文件（OCR 文字识别）。"""
    from core.rag.document.ocr import ocr_image

    text = ocr_image(path)
    return LoadedDocument(
        content=text,
        source_path=path,
        file_type="image",
        metadata={"file_name": os.path.basename(path)},
    )


# ========== 统一加载器 ==========

# 支持的扩展名映射
_LOADER_MAP = {
    ".txt": _load_txt,
    ".md": _load_txt,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
    ".png": _load_image,
    ".jpg": _load_image,
    ".jpeg": _load_image,
    ".bmp": _load_image,
    ".tiff": _load_image,
}


def load_document(path: str) -> LoadedDocument:
    """加载单个文档，自动识别格式。

    Args:
        path: 文件路径

    Returns:
        LoadedDocument

    Raises:
        ValueError: 不支持的格式
        ImportError: 缺少对应依赖
    """
    ext = Path(path).suffix.lower()
    if ext not in _LOADER_MAP:
        supported = ", ".join(_LOADER_MAP.keys())
        raise ValueError(f"不支持的文件格式: {ext}。支持: {supported}")

    loader = _LOADER_MAP[ext]
    logger.info("加载文档: %s (type=%s)", path, ext)
    return loader(path)


def load_documents(
    paths: Union[str, List[str]],
    recursive: bool = False,
) -> List[LoadedDocument]:
    """批量加载文档。

    Args:
        paths: 文件路径 或 目录路径
        recursive: 是否递归子目录

    Returns:
        LoadedDocument 列表
    """
    if isinstance(paths, str):
        paths = [paths]

    file_paths = []
    for p in paths:
        pp = Path(p)
        if pp.is_file():
            file_paths.append(str(pp))
        elif pp.is_dir():
            glob = pp.rglob("*") if recursive else pp.glob("*")
            file_paths.extend(
                str(f) for f in glob
                if f.is_file() and f.suffix.lower() in _LOADER_MAP
            )

    documents = []
    for fp in file_paths:
        try:
            doc = load_document(fp)
            if doc.content.strip():
                documents.append(doc)
        except Exception as e:
            logger.warning("加载文档失败 %s: %s", fp, e)

    logger.info("批量加载完成: %d/%d 个文档", len(documents), len(file_paths))
    return documents


def load_and_chunk(
    paths: Union[str, List[str]],
    recursive: bool = False,
    parent_chunk_size: int = 800,
    child_chunk_size: int = 200,
) -> tuple:
    """加载文档并执行父子分块。

    一站式处理：加载 → 清洗 → 分块 → 返回可用于建索引的数据。

    Returns:
        (parents, children, child_to_parent_map)
    """
    from core.rag.chunker import ParentChildChunker, clean_text

    documents = load_documents(paths, recursive)
    chunker = ParentChildChunker(
        parent_chunk_size=parent_chunk_size,
        child_chunk_size=child_chunk_size,
    )

    doc_inputs = [(clean_text(doc.content), doc.metadata) for doc in documents]
    return chunker.chunk_documents(doc_inputs)